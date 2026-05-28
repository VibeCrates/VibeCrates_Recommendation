"""
Crates_4.pptx를 베이스로 내용만 교체해 Crates_5.pptx 생성.
기존 슬라이드를 직접 수정하여 Duplicate name 문제 회피.
"""

import copy
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

TEMPLATE = "Crates_4.pptx"
OUTPUT   = "Crates_5.pptx"

shutil.copy2(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)
slides = list(prs.slides)

PURPLE = RGBColor(0x7C, 0x3A, 0xED)
DARK   = RGBColor(0x1F, 0x29, 0x37)
GRAY   = RGBColor(0x53, 0x53, 0x53)
GREEN  = RGBColor(0x05, 0x96, 0x69)
RED    = RGBColor(0xDC, 0x26, 0x26)


# ── 헬퍼 ────────────────────────────────────────────────────

def clear_txbody(txBody):
    """txBody 안의 <a:p> 전부 제거."""
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)


def make_para(lines):
    """
    lines: str 또는 [(text, bold, color_or_None), ...] 리스트
    → <a:p> XML Element 반환
    """
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    p = etree.SubElement(etree.Element('dummy'), f'{{{NS}}}p')
    p = etree.Element(f'{{{NS}}}p')

    pPr = etree.SubElement(p, f'{{{NS}}}pPr')
    pPr.set('algn', 'l')
    spcBef = etree.SubElement(pPr, f'{{{NS}}}spcBef')
    etree.SubElement(spcBef, f'{{{NS}}}spcPts').set('val', '0')
    spcAft = etree.SubElement(pPr, f'{{{NS}}}spcAft')
    etree.SubElement(spcAft, f'{{{NS}}}spcPts').set('val', '0')

    if isinstance(lines, str):
        segments = [(lines, False, None)]
    elif isinstance(lines, list):
        # 각 항목을 (text, bold, color) 튜플로 정규화
        segments = []
        for seg in lines:
            if isinstance(seg, str):
                segments.append((seg, False, None))
            elif len(seg) == 1:
                segments.append((seg[0], False, None))
            elif len(seg) == 2:
                segments.append((seg[0], seg[1], None))
            else:
                segments.append((seg[0], seg[1], seg[2]))
    else:
        segments = [(str(lines), False, None)]

    for text, bold, color in segments:
        r = etree.SubElement(p, f'{{{NS}}}r')
        rPr = etree.SubElement(r, f'{{{NS}}}rPr')
        rPr.set('lang', 'ko')
        if bold:
            rPr.set('b', '1')
        if color:
            solidFill = etree.SubElement(rPr, f'{{{NS}}}solidFill')
            srgbClr = etree.SubElement(solidFill, f'{{{NS}}}srgbClr')
            srgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')
        t = etree.SubElement(r, f'{{{NS}}}t')
        t.text = text

    return p


def set_slide(slide, title_str, body_items):
    """슬라이드의 title / body 텍스트를 교체."""
    shapes = list(slide.shapes)
    text_shapes = [s for s in shapes if s.has_text_frame]
    if len(text_shapes) < 1:
        return

    # title
    title_shape = text_shapes[0]
    txBody = title_shape.text_frame._txBody
    clear_txbody(txBody)
    txBody.append(make_para(title_str))

    # body
    if len(text_shapes) < 2 or not body_items:
        return
    body_shape = text_shapes[1]
    txBody2 = body_shape.text_frame._txBody
    clear_txbody(txBody2)
    for item in body_items:
        txBody2.append(make_para(item))


def add_slide_like(prs, src_slide):
    """src_slide 레이아웃으로 새 슬라이드 추가."""
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # 도형 초기화
    sp_tree = new_slide.shapes._spTree
    for el in list(sp_tree):
        tag = el.tag.split('}')[-1]
        if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
            sp_tree.remove(el)
    # src 도형 복사
    for shape in src_slide.shapes:
        sp_tree.append(copy.deepcopy(shape._element))
    return new_slide


# ── 슬라이드 내용 정의 ────────────────────────────────────────

# (title, body_items)
# body_items 항목: str  또는  [(text, bold, color), ...]

CONTENT = [
    # 슬라이드 1 — 타이틀 (내용 그대로 유지, 제목만 변경)
    (
        "VibeCrates Recommendation System",
        [("데이터 전처리 & 모델 설계 현황", False, GRAY)],
    ),
    # 슬라이드 2 — 프로젝트 개요
    (
        "프로젝트 개요",
        [
            [("VibeCrates", True, PURPLE), (" — 영화·음악·도서 멀티모달 콘텐츠 추천 시스템", False, None)],
            "",
            [("핵심 구조", True, DARK)],
            "  이미지 + 텍스트 + 사용자 쿼리 → 768차원 벡터 → ANN 검색",
            "",
            [("모델", True, DARK)],
            "  DualEncoderModel  ( TextBlock / ImageBlock / QueryBlock / ContentBlock )",
            "",
            [("학습", True, DARK)],
            "  Stage 1 : InfoNCE Loss   |   Stage 2 : KL Divergence",
            "",
            [("데이터 규모", True, DARK)],
            "  음악 40,036개  |  도서 40,000개  |  영화 40,108개",
        ],
    ),
    # 슬라이드 3 — 음악 데이터 구성
    (
        "음악 데이터 구성",
        [
            [("데이터 소스", True, PURPLE)],
            "  • Music.csv    — 이미지(Spotify CDN) + audio features  (65,795개)",
            "  • songs.csv   — 가사 + 장르  (24,029,011개)",
            "",
            [("Intersection  ( Spotify ID 기준 )", True, DARK)],
            "  두 파일의 교집합 → 이미지 & 가사 모두 존재",
            "",
            [("Instrumental 보완", True, DARK)],
            "  Music.csv에서 instrumentalness ≥ 0.5 인 항목 별도 추출·추가",
            [("  → 가사 없음 4,092개 — Text 처리 필요", False, RED)],
            "",
            [("최종", True, DARK)],
            [("  총 40,036개", True, None), ("   ( 일반 35,944 + Instrumental 4,092 )", False, GRAY)],
        ],
    ),
    # 슬라이드 4 — 음악 이미지 유효성
    (
        "음악 이미지 유효성 검증",
        [
            [("현황", True, PURPLE)],
            "  img 컬럼 40,036개  →  URL 보유 21,742개  /  'no' 18,292개  /  null 2개",
            "",
            [("Spotify API 시도 → 실패", True, RED)],
            "  배치 엔드포인트 (/v1/tracks?ids=…) : 403 Forbidden",
            "  단일 요청 대량 시도 → Retry-After 82,970초 (≈ 23시간) rate limit",
            "",
            [("iTunes Search API 전환", True, GREEN)],
            "  API 키 불필요  /  아티스트명 + 트랙명으로 검색",
            "  artworkUrl100 → 600×600 업스케일",
            "  초당 10req, 18,244개 처리 중 (약 30분 소요)",
        ],
    ),
    # 슬라이드 5 — 도서 데이터 구성
    (
        "도서 데이터 구성",
        [
            [("원본", True, PURPLE), ("  Books.csv — 271,360개", False, None)],
            "",
            [("Step 1  이미지 유효성 필터링", True, DARK)],
            "  Amazon CDN — 200 OK 이지만 1×1 GIF / 160×107 JPEG 플레이스홀더 존재",
            "  GET 스트리밍 + MD5 체크 방식으로 무효 이미지 제거 (ThreadPool 200 workers)",
            [("  → 149,737개", True, GREEN)],
            "",
            [("Step 2  Open Library API로 카테고리 추가", True, DARK)],
            "  ISBN 배치 요청, subjects 첫 항목을 category로 저장",
            [("  → 123,774개  ( category / main_category / raw_category )", True, GREEN)],
            "",
            [("Step 3  카테고리 정규화", True, DARK)],
            [("  고유 카테고리 15,543개 → 25개로 통합", False, None)],
            "  Fiction / Mystery·Thriller / Children's·YA / Romance …",
        ],
    ),
    # 슬라이드 6 — 도서 Popularity & 샘플링
    (
        "도서 Popularity 정의 & 샘플링",
        [
            [("Popularity Score  ( Open Library Search API )", True, PURPLE)],
            "  popularity_score = already_read_count + want_to_read_count + edition_count × 10",
            "  already_read·want_to_read: 실제 독자 행동 / edition_count × 10: 역사적 인기 보완",
            "",
            [("분포", True, DARK)],
            "  score > 0 :   8,884개  ← 전체 포함 확정",
            "  score = 0 : 114,890개  ← 층화 샘플링 대상",
            "",
            [("층화 샘플링  ( 목표 40,000개 )", True, DARK)],
            "  score > 0 고정 8,884개 + score = 0에서 main_category 비율 유지 31,116개",
            [("  → Books_final.csv  40,000개 완료", True, GREEN)],
            "",
            "  Other 19,355  |  Fiction 9,463  |  Children's/YA 1,336  |  …",
        ],
    ),
    # 슬라이드 7 — 모델 아키텍처
    (
        "모델 아키텍처 — DualEncoderModel",
        [
            [("인코더 블록", True, PURPLE)],
            "  TextBlock  :  SBERT (paraphrase-multilingual-mpnet-base-v2) + LoRA(r=16) + MLP → 768dim",
            "  ImageBlock :  CLIP Vision Encoder (Bingsu/clip-vit-large-patch14-ko) + MLP → 768dim",
            "  QueryBlock :  CLIP Text Encoder (동일 모델) + MLP → 768dim",
            "",
            [("ContentBlock", True, PURPLE)],
            "  concat(z_image, z_text)  [1536dim]",
            "  → Linear → LayerNorm → GELU → Dropout  (×3)",
            "  → L2 normalize → 768dim",
            "",
            [("학습 전략", True, DARK)],
            "  사전학습 파라미터 전부 Freeze — LoRA 어댑터 + MLP 헤드만 학습",
        ],
    ),
    # 슬라이드 8 — 학습 파이프라인
    (
        "학습 파이프라인 — TwoStageTrainer",
        [
            [("Stage 1  — LoRA + MLP 학습  ( 10 에폭, lr=1e-4 )", True, PURPLE)],
            "  InfoNCELoss  — 대칭적 양방향 Cross-Entropy, temperature=0.07",
            "  페어 3쌍 : 텍스트-쿼리 / 텍스트-이미지 / 이미지-쿼리  합산",
            "",
            [("Stage 2  — ContentBlock 학습  ( 15 에폭, lr=5e-5 )", True, PURPLE)],
            "  KLDivergenceLoss  — 지식 증류 (log_target=True)",
            "  Teacher: z_query (no_grad)  vs  Student: z_content",
            "",
            [("공통 설정", True, DARK)],
            "  batch_size=32  /  early_stopping_patience=3",
            "  LoRA: r=16, alpha=32, dropout=0.05",
            "",
            [("TrainingHistory", True, DARK)],
            "  배치·에폭 단위 메트릭 기록 + JSON 저장/로드",
        ],
    ),
    # 슬라이드 9 — 데이터 처리 모듈
    (
        "데이터 처리 모듈",
        [
            [("완성  ✓", True, GREEN)],
            "  MultiModalDataset  — (content_text, image_path, query) 3쌍 반환",
            "                          PIL 224×224 리사이즈  +  collate_fn",
            "  DataLoader         — CSV → Dataset → train/val/test split",
            "                          pin_memory=True, train만 shuffle",
            "",
            [("미완성  ✗", True, RED)],
            "  DataPreprocessor  — 골격만 존재, 아래 전부 TODO",
            "    • 결측치 처리",
            "    • 범주형 인코딩",
            "    • 피처 정규화",
            "    • 사용자 프로필 생성",
            "    • 아이템 피처 생성",
        ],
    ),
    # 슬라이드 10 — 전체 구현 현황
    (
        "전체 구현 현황",
        [
            [("완성  ✓", True, GREEN)],
            "  models/base.py · recommender.py · utils.py(일부)",
            "  training/config.py · losses.py · trainer.py · history.py",
            "  data/dataset.py · loader.py",
            "",
            [("미완성  ✗", True, RED)],
            "  data/preprocessing.py  — DataPreprocessor (TODO)",
            "  src/api/               — routes, dependencies 미구현",
            "  scripts/evaluate.py    — 평가 스크립트",
            "  scripts/inference.py   — 추론 스크립트",
            "",
            [("데이터 전처리 진행중", True, DARK)],
            "  음악 이미지 수집 (iTunes API)  — 진행 중",
            "  Instrumental 텍스트 생성       — 방법 검토 중",
            "  도서 Books_final.csv           — 완료",
        ],
    ),
    # 슬라이드 11 — Q&A
    (
        "Q & A",
        [
            "발표를 들어주셔서 감사합니다.",
            "",
            [("Crates Team", True, PURPLE)],
        ],
    ),
]

# ── 슬라이드 수 맞추기 ────────────────────────────────────────
# 원본 12개 → 목표 11개: 마지막 슬라이드 삭제 후 내용 교체

# 원본 슬라이드가 12개, 목표 11개 → 12번째 슬라이드 XML 제거
sldIdLst = prs._element.find(qn('p:sldIdLst'))
sldIds = list(sldIdLst)

# 슬라이드가 목표보다 많으면 뒤에서 제거
while len(sldIds) > len(CONTENT):
    sldIdLst.remove(sldIds[-1])
    sldIds = list(sldIdLst)

# 슬라이드가 목표보다 적으면 마지막 슬라이드 복제해서 추가
src_for_add = slides[-2]  # Q&A 전 슬라이드 스타일
while len(list(sldIdLst)) < len(CONTENT):
    add_slide_like(prs, src_for_add)

# 저장 후 다시 열기 (슬라이드 목록 갱신)
prs.save(OUTPUT)
prs = Presentation(OUTPUT)
slides = list(prs.slides)

# ── 내용 채우기 ───────────────────────────────────────────────
for i, (title, body) in enumerate(CONTENT):
    set_slide(slides[i], title, body)

prs.save(OUTPUT)
print(f"저장 완료: {OUTPUT}  ({len(slides)}개 슬라이드)")
